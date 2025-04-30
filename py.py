import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# --- Presentation Content (Based on your previous structure) ---
# Structure: (layout_type, title, list_of_bullets)
# layout_type: 'title_slide', 'content_slide', 'blank_slide'

presentation_content = [
    (
        "title_slide",
        "Dezvoltarea unei Aplicații Web Interactive Tic Tac Toe",
        [
            "Raportul Stagiului de Practică",
            "Student: Grigoriev Artur, Grupa PTPP-241",
            "Specializarea: Programarea și testarea produselor program",
            "Instituția: Colegiul Universității Tehnice a Moldovei",
            "Conducător: Moraru Magdalena",
            "Chișinău 2025",
        ],
    ),
    (
        "content_slide",
        "Introducere & Context",
        [
            "Scop Proiect: Dezvoltarea unei aplicații web Tic Tac Toe interactive și educative.",
            "Context: Proiect de stagiu pentru specializarea „Programarea și Testarea Produselor Program”.",
            "Motivație:",
            "  - Aplicarea cunoștințelor teoretice (programare, modelare date, AI).",
            "  - Explorarea tehnologiilor web și dezvoltării jocurilor.",
            "  - Crearea unei experiențe utilizator atractive.",
            "Scop Raport: Acoperă analiza domeniului, designul sistemului, implementarea și testarea (22 pagini).",
        ],
    ),
    (
        "content_slide",
        "Obiective & Cerințe Proiect",
        [
            "Scop Principal: Crearea unei experiențe digitale Tic Tac Toe atractive și interactive.",
            "Obiective Cheie:",
            "  - Platformă robustă și intuitivă (web).",
            "  - Moduri: Single-player (vs. AI) și Multiplayer Local.",
            "  - Configurare: Dimensiune grilă (3x3, 5x5), Simbol jucător (X/O).",
            "  - Implementare AI competitiv (algoritmi strategici).",
            "  - Autentificare utilizator și gestionare sesiune.",
            "  - Afișare scoruri și feedback vizual clar (animații).",
            "  - Design responsiv și modular.",
        ],
    ),
    (
        "content_slide",
        "Tehnologii & Arhitectură",
        [
            "Frontend: HTML5, CSS3, JavaScript (ES6+)",
            "Motor Joc: Phaser.js [1] - Grafică 2D, animații, input, scene web.",
            "Autentificare: Mecanism simplu folosind localStorage.",
            "Arhitectură: Modulară",
            "  - Logică autentificare",
            "  - Gestionare stare joc",
            "  - Randare UI (via Phaser)",
            "  - Logică AI",
        ],
    ),
    (
        "content_slide",
        "Implementare - Logică & UI",
        [
            "Autentificare: Verificare localStorage, interfață Login/Register (Fig 2.1, 2.2), Logout.",
            "Setup Joc (Phaser):",
            "  - preload(): Încărcare resurse (imagini X/O - Anexa A).",
            "  - create(): Desenare grilă dinamică, celule interactive (Pag. 11).",
            "Logică Joc (TakeTurn, Check):",
            "  - Gestionare click, validare mutări.",
            "  - Actualizare 'matrix' internă.",
            "  - Plasare sprite X/O cu animație.",
            "  - Alternare ture.",
            "  - Check() (Anexa D): Verificare condiții câștig/remiză.",
            "  - Actualizare și afișare scoruri (displayWinner).",
        ],
    ),
    (
        "content_slide",
        "Implementare - Adversar AI",
        [
            "Scop: Creare AI competitiv pentru modul single-player.",
            "Strategie Hibridă:",
            "  1. Verificări Imediate (findImmediate): Câștig/Blocare în 1 mutare.",
            "  2. Algoritm Minimax [5]: Explorare arbore joc (dacă nu există mutare imediată).",
            "  3. Optimizare Alpha-Beta: Eliminare ramuri irelevante.",
            "  4. Funcție Euristică (heuristic): Evaluare stări non-terminale.",
            "  5. Limite Adâncime/Noduri (MM_DEPTH, MM_NODES): Asigurare timp răspuns rezonabil.",
        ],
    ),
    (
        "content_slide",
        "Interfață Utilizator & Testare",
        [
            "Ecrane Cheie:",
            "  - Login/Înregistrare (Fig 2.1, 2.2).",
            "  - Meniu Principal: Selectare simbol, adversar, dimensiune tablă (Fig 2.3).",
            "  - Ecran Joc: Tablă centrală, scor, elemente vizuale clare (Fig 2.4).",
            "  - Navigare: Butoane Meniu/Logout (Fig 2.5).",
            "Testare: Focus pe corectitudinea funcțională a modulelor (autentificare, logică joc, AI, UI).",
        ],
    ),
    (
        "content_slide",
        "Concluzii & Învățăminte",
        [
            "Realizări:",
            "  - Aplicație Tic Tac Toe web funcțională și interactivă.",
            "  - Implementare obiective: moduri multiple, grilă configurabilă, scor, autentificare.",
            "  - Integrare AI competitiv (Minimax + Alpha-Beta).",
            "  - Utilizare eficientă Phaser.js.",
            "Învățăminte:",
            "  - Aplicare practică JS, HTML, CSS.",
            "  - Experiență cu motorul Phaser.js.",
            "  - Implementare algoritmi căutare AI.",
            "  - Înțelegere structură modulară și ciclu viață aplicație web.",
            "  - Experiență valoroasă în transpunerea cerințelor în soluții software.",
        ],
    ),
    ("blank_slide", "Întrebări?", ["Mulțumesc!"]),
]

# --- Formatting Settings (Based on Requirements) ---
FONT_NAME = "Times New Roman"
FONT_COLOR_RGB = RGBColor(0, 0, 0)
TITLE_FONT_SIZE = Pt(28)  # Adjusted for PPT title visibility
SUBTITLE_FONT_SIZE = Pt(18) # Adjusted for PPT subtitle visibility
BODY_FONT_SIZE = Pt(14)  # Slightly larger than 12pt for readability on screen
CHAPTER_TITLE_FONT_SIZE = Pt(24) # Adjusted for PPT content slide title

# --- Helper Function to Apply Formatting ---
def set_run_format(run, size=BODY_FONT_SIZE, bold=False):
    """Applies font name, size, bold, and color to a text run."""
    font = run.font
    font.name = FONT_NAME
    font.size = size
    font.bold = bold
    font.color.rgb = FONT_COLOR_RGB

# --- Create Presentation ---
prs = Presentation()

# --- Define Layouts (Indices might vary slightly based on default template) ---
title_slide_layout = prs.slide_layouts[0]  # Usually Title Slide
content_slide_layout = prs.slide_layouts[1]  # Usually Title and Content
blank_slide_layout = prs.slide_layouts[5]  # Usually Blank

# --- Add Slides and Content ---
for index, (layout_type, title_text, bullets) in enumerate(
    presentation_content
):
    if layout_type == "title_slide":
        slide = prs.slides.add_slide(title_slide_layout)
        # Title
        title_shape = slide.shapes.title
        title_shape.text = title_text
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        set_run_format(
            title_shape.text_frame.paragraphs[0].runs[0],
            size=TITLE_FONT_SIZE,
            bold=True,
        )

        # Subtitle (using the body placeholder for multiple lines)
        subtitle_shape = slide.placeholders[1]
        tf = subtitle_shape.text_frame
        tf.clear() # Clear existing text
        tf.word_wrap = True
        for i, bullet_text in enumerate(bullets):
            p = tf.add_paragraph()
            p.text = bullet_text
            p.alignment = PP_ALIGN.CENTER
            # Apply formatting to the first run (assuming simple text)
            if p.runs:
                set_run_format(p.runs[0], size=SUBTITLE_FONT_SIZE, bold=(i==0)) # Bold first subtitle line maybe

    elif layout_type == "content_slide":
        slide = prs.slides.add_slide(content_slide_layout)
        # Title
        title_shape = slide.shapes.title
        title_shape.text = title_text
        # Apply formatting to the first run
        if title_shape.text_frame.paragraphs[0].runs:
             set_run_format(
                title_shape.text_frame.paragraphs[0].runs[0],
                size=CHAPTER_TITLE_FONT_SIZE,
                bold=True,
            )
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT # Chapter titles left aligned usually

        # Body Content (Bullets)
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear() # Clear default text
        tf.word_wrap = True

        for bullet_text in bullets:
            p = tf.add_paragraph()
            # Handle indentation based on leading spaces
            indent_level = 0
            stripped_text = bullet_text.lstrip()
            spaces = len(bullet_text) - len(stripped_text)
            if spaces >= 2: # Simple indentation check
                indent_level = 1
            if spaces >= 4:
                 indent_level = 2

            p.text = stripped_text
            p.level = indent_level
            # Justify is possible but often looks odd in PPT, Left is safer
            p.alignment = PP_ALIGN.LEFT
            # Apply formatting to the first run
            if p.runs:
                set_run_format(p.runs[0], size=BODY_FONT_SIZE)

    elif layout_type == "blank_slide":
        slide = prs.slides.add_slide(blank_slide_layout)
        # Add a text box for the title manually if needed on blank slide
        # For simplicity, let's add title and one line of text
        left = top = width = height = Inches(1.0) # Example position/size
        txBox = slide.shapes.add_textbox(left, top, Inches(8.5), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True

        # Title for Q&A
        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.alignment = PP_ALIGN.CENTER
        if p_title.runs:
            set_run_format(p_title.runs[0], size=TITLE_FONT_SIZE, bold=True)

        # Additional text
        for bullet_text in bullets:
             p_text = tf.add_paragraph()
             p_text.text = bullet_text
             p_text.alignment = PP_ALIGN.CENTER
             if p_text.runs:
                 set_run_format(p_text.runs[0], size=SUBTITLE_FONT_SIZE)


# --- Save the Presentation ---
output_filename = "Prezentare_TicTacToe_Artur_Grigoriev.pptx"
prs.save(output_filename)

print(f"Prezentarea a fost generată: {output_filename}")
